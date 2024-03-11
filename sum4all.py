import requests
import json
import re
import plugins
from bridge.reply import Reply, ReplyType
from bridge.context import ContextType
from channel.chat_message import ChatMessage
from plugins import *
from common.log import logger
from common.expired_dict import ExpiredDict
import os
from docx import Document
import markdown
import tiktoken
import jieba
import fitz
from openpyxl import load_workbook
import csv
from bs4 import BeautifulSoup
from pptx import Presentation
from PIL import Image
import base64


EXTENSION_TO_TYPE = {
    'pdf': 'pdf',
    'doc': 'docx', 'docx': 'docx',
    'md': 'md',
    'txt': 'txt',
    'xls': 'excel', 'xlsx': 'excel',
    'csv': 'csv',
    'html': 'html', 'htm': 'html',
    'ppt': 'ppt', 'pptx': 'ppt'
}

@plugins.register(
    name="sum4all",
    desire_priority=2,
    desc="A plugin for summarizing all things",
    version="0.7.3",
    author="fatwang2",
)

class sum4all(Plugin):
    def __init__(self):
        super().__init__()
        try:
            curdir = os.path.dirname(__file__)
            config_path = os.path.join(curdir, "config.json")
            if os.path.exists(config_path):
                with open(config_path, "r", encoding="utf-8") as f:
                    self.config = json.load(f)
            else:
                # 使用父类的方法来加载配置
                self.config = super().load_config()

                if not self.config:
                    raise Exception("config.json not found")
            # 设置事件处理函数
            self.handlers[Event.ON_HANDLE_CONTEXT] = self.on_handle_context
            self.params_cache = ExpiredDict(300)

            
            # 从配置中提取所需的设置
            self.keys = self.config.get("keys", {})
            self.url_sum = self.config.get("url_sum", {})
            self.search_sum = self.config.get("search_sum", {})
            self.file_sum = self.config.get("file_sum", {})
            self.image_sum = self.config.get("image_sum", {})
            self.note = self.config.get("note", {})

            self.sum4all_key = self.keys.get("sum4all_key", "")
            self.gemini_key = self.keys.get("gemini_key", "")
            self.bibigpt_key = self.keys.get("bibigpt_key", "")
            self.outputLanguage = self.keys.get("outputLanguage", "zh-CN")
            self.opensum_key = self.keys.get("opensum_key", "")
            self.open_ai_api_key = self.keys.get("open_ai_api_key", "")
            self.model = self.keys.get("model", "gpt-3.5-turbo")
            self.open_ai_api_base = self.keys.get("open_ai_api_base", "https://api.openai.com/v1")
            self.xunfei_app_id = self.keys.get("xunfei_app_id", "")
            self.xunfei_api_key = self.keys.get("xunfei_api_key", "")
            self.xunfei_api_secret = self.keys.get("xunfei_api_secret", "")
            self.perplexity_key = self.keys.get("perplexity_key", "")
            self.flomo_key = self.keys.get("flomo_key", "")

            # 提取sum服务的配置
            self.url_sum_enabled = self.url_sum.get("enabled", False)
            self.url_sum_service = self.url_sum.get("service", "")
            self.url_sum_group = self.url_sum.get("group", True)
            self.url_sum_qa_prefix = self.url_sum.get("qa_prefix", "问")
            self.url_sum_prompt = self.url_sum.get("prompt", "")

            self.search_sum_enabled = self.search_sum.get("enabled", False)
            self.search_sum_service = self.search_sum.get("service", "")
            self.search_sum_group = self.search_sum.get("group", True)
            self.search_sum_search_prefix = self.search_sum.get("search_prefix", "搜")
            self.search_sum_prompt = self.search_sum.get("prompt", "")

            self.file_sum_enabled = self.file_sum.get("enabled", False)
            self.file_sum_service = self.file_sum.get("service", "")
            self.file_sum_group = self.file_sum.get("group", True)
            self.file_sum_qa_prefix = self.file_sum.get("qa_prefix", "问")
            self.file_sum_prompt = self.file_sum.get("prompt", "")

            self.image_sum_enabled = self.image_sum.get("enabled", False)
            self.image_sum_service = self.image_sum.get("service", "")
            self.image_sum_group = self.image_sum.get("group", True)
            self.image_sum_qa_prefix = self.image_sum.get("qa_prefix", "问")
            self.image_sum_prompt = self.image_sum.get("prompt", "")

            self.note_enabled = self.note.get("enabled", False)
            self.note_service = self.note.get("service", "")
            self.note_prefix = self.note.get("prefix", "记")

            # 初始化成功日志
            logger.info("[sum4all] inited.")
        except Exception as e:
            # 初始化失败日志
            logger.warn(f"sum4all init failed: {e}")
    def on_handle_context(self, e_context: EventContext):
        context = e_context["context"]
        if context.type not in [ContextType.TEXT, ContextType.SHARING,ContextType.FILE,ContextType.IMAGE]:
            return
        msg: ChatMessage = e_context["context"]["msg"]
        user_id = msg.from_user_id
        content = context.content
        isgroup = e_context["context"].get("isgroup", False)

        url_match = re.match('https?://(?:[-\w.]|(?:%[\da-fA-F]{2}))+', content)
        unsupported_urls = re.search(r'.*finder\.video\.qq\.com.*|.*support\.weixin\.qq\.com/update.*|.*support\.weixin\.qq\.com/security.*|.*mp\.weixin\.qq\.com/mp/waerrpage.*', content)

            # 检查输入是否以"搜索前缀词" 开头
        if content.startswith(self.search_sum_search_prefix) and self.search_sum_enabled:
            # 如果消息来自一个群聊，并且你不希望在群聊中启用搜索功能，直接返回
            if isgroup and not self.search_sum_group:
                return
            # Call new function to handle search operation
            self.call_service(content, e_context, "search")
            return
        
        if user_id in self.params_cache and ('last_file_content' in self.params_cache[user_id] or 'last_image_base64' in self.params_cache[user_id] or 'last_url' in self.params_cache[user_id]):
            # 如果存在最近一次处理的文件路径，触发文件理解函数
            if 'last_file_content' in self.params_cache[user_id] and content.startswith(self.file_sum_qa_prefix):
                logger.info('Content starts with the file_sum_qa_prefix.')
                # 去除关键词和紧随其后的空格
                new_content = content[len(self.file_sum_qa_prefix):]
                self.params_cache[user_id]['prompt'] = new_content
                logger.info('params_cache for user has been successfully updated.')            
                self.handle_file(self.params_cache[user_id]['last_file_content'], e_context)
            # 如果存在最近一次处理的图片路径，触发图片理解函数
            elif 'last_image_base64' in self.params_cache[user_id] and content.startswith(self.image_sum_qa_prefix):
                logger.info('Content starts with the image_sum_qa_prefix.')
                # 去除关键词和紧随其后的空格
                new_content = content[len(self.image_sum_qa_prefix):]
                self.params_cache[user_id]['prompt'] = new_content
                logger.info('params_cache for user has been successfully updated.')            
                if self.image_sum_service == "xunfei":
                    self.handle_xunfei_image(self.params_cache[user_id]['last_image_base64'], e_context)
                elif self.image_sum_service == "openai":
                    self.handle_openai_image(self.params_cache[user_id]['last_image_base64'], e_context)
                elif self.image_sum_service == "gemini":
                    self.handle_gemini_image(self.params_cache[user_id]['last_image_base64'], e_context)
            # 如果存在最近一次处理的URL，触发URL理解函数
            elif 'last_url' in self.params_cache[user_id] and content.startswith(self.url_sum_qa_prefix):
                logger.info('Content starts with the url_sum_qa_prefix.')
                # 去除关键词和紧随其后的空格
                new_content = content[len(self.url_sum_qa_prefix):]
                self.params_cache[user_id]['prompt'] = new_content
                logger.info('params_cache for user has been successfully updated.')            
                self.call_service(self.params_cache[user_id]['last_url'], e_context ,"sum")
            elif 'last_url' in self.params_cache[user_id] and content.startswith(self.note_prefix) and self.note_enabled and not isgroup:
                logger.info('Content starts with the note_prefix.')
                new_content = content[len(self.note_prefix):]
                self.params_cache[user_id]['note'] = new_content
                logger.info('params_cache for user has been successfully updated.')  
                self.call_service(self.params_cache[user_id]['last_url'], e_context, "note")
        if context.type == ContextType.FILE:
            if isgroup and not self.file_sum_group:
                # 群聊中忽略处理文件
                logger.info("群聊消息，文件处理功能已禁用")
                return
            logger.info("on_handle_context: 处理上下文开始")
            context.get("msg").prepare()
            file_path = context.content
            logger.info(f"on_handle_context: 获取到文件路径 {file_path}")
            
            
            # 检查是否应该进行文件总结
            if self.file_sum_enabled:
                # 更新params_cache中的last_file_content
                self.params_cache[user_id] = {}
                file_content = self.extract_content(file_path)
                if file_content is None:
                    logger.info("文件内容无法提取，跳过处理")
                else:
                    self.params_cache[user_id]['last_file_content'] = file_content
                    logger.info('Updated last_file_content in params_cache for user.')
                    self.handle_file(file_content, e_context)
            else:
                logger.info("文件总结功能已禁用，不对文件内容进行处理")
            # 删除文件
            os.remove(file_path)
            logger.info(f"文件 {file_path} 已删除")
        elif context.type == ContextType.IMAGE:
            if isgroup and not self.image_sum_group:
                # 群聊中忽略处理图片
                logger.info("群聊消息，图片处理功能已禁用")
                return
            logger.info("on_handle_context: 开始处理图片")
            context.get("msg").prepare()
            image_path = context.content
            logger.info(f"on_handle_context: 获取到图片路径 {image_path}")
            
            
            # 检查是否应该进行图片总结
            if self.image_sum_enabled:
                # 将图片路径转换为Base64编码的字符串
                base64_image = self.encode_image_to_base64(image_path)
                # 更新params_cache中的last_image_path
                self.params_cache[user_id] = {}
                self.params_cache[user_id]['last_image_base64'] = base64_image
                logger.info('Updated last_image_base64 in params_cache for user.')
                if self.image_sum_service == "xunfei":
                    self.handle_xunfei_image(base64_image, e_context)
                elif self.image_sum_service == "openai":
                    self.handle_openai_image(base64_image, e_context)
                elif self.image_sum_service == "gemini":
                    self.handle_gemini_image(base64_image, e_context)
            else:
                logger.info("图片总结功能已禁用，不对图片内容进行处理")
            # 删除文件
            os.remove(image_path)
            logger.info(f"文件 {image_path} 已删除")
        elif context.type == ContextType.SHARING:  #匹配卡片分享
            if unsupported_urls:  #匹配不支持总结的卡片
                if isgroup:  ##群聊中忽略
                    return
                else:  ##私聊回复不支持
                    logger.info("[sum4all] Unsupported URL : %s", content)
                    reply = Reply(type=ReplyType.TEXT, content="不支持总结小程序和视频号")
                    e_context["reply"] = reply
                    e_context.action = EventAction.BREAK_PASS
            else:  #匹配支持总结的卡片
                if isgroup:  #处理群聊总结
                    if self.url_sum_group:  #group_sharing = True进行总结，False则忽略。
                        logger.info("[sum4all] Summary URL : %s", content)
                        # 更新params_cache中的last_url
                        self.params_cache[user_id] = {}
                        self.params_cache[user_id]['last_url'] = content
                        logger.info('Updated last_url in params_cache for user.')
                        self.call_service(content, e_context, "sum")
                        return
                    else:
                        return
                else:  #处理私聊总结
                    logger.info("[sum4all] Summary URL : %s", content)
                    # 更新params_cache中的last_url
                    self.params_cache[user_id] = {}
                    self.params_cache[user_id]['last_url'] = content
                    logger.info('Updated last_url in params_cache for user.')
                    self.call_service(content, e_context, "sum")
                    return
            
        elif url_match: #匹配URL链接
            if unsupported_urls:  #匹配不支持总结的网址
                logger.info("[sum4all] Unsupported URL : %s", content)
                reply = Reply(type=ReplyType.TEXT, content="不支持总结小程序和视频号")
                e_context["reply"] = reply
                e_context.action = EventAction.BREAK_PASS
            else:
                logger.info("[sum4all] Summary URL : %s", content)
                # 更新params_cache中的last_url
                self.params_cache[user_id] = {}
                self.params_cache[user_id]['last_url'] = content
                logger.info('Updated last_url in params_cache for user.')
                self.call_service(content, e_context, "sum")
                return
    def call_service(self, content, e_context, service_type):
        if service_type == "search":
            if self.search_sum_service == "openai" or self.search_sum_service == "sum4all" or self.search_sum_service == "gemini":
                self.handle_search(content, e_context)
            elif self.search_sum_service == "perplexity":
                self.handle_perplexity(content, e_context)
        elif service_type == "sum":
            if self.url_sum_service == "bibigpt":
                self.handle_bibigpt(content, e_context)
            elif self.url_sum_service == "openai" or self.url_sum_service == "sum4all" or self.url_sum_service == "gemini":
                self.handle_url(content, e_context)
            elif self.url_sum_service == "opensum":
                self.handle_opensum(content, e_context)
        elif service_type == "note":
            if self.note_service == "flomo":
                self.handle_note(content, e_context)
    def handle_note(self,link,e_context):
        msg: ChatMessage = e_context["context"]["msg"]
        user_id = msg.from_user_id
        title = self.params_cache[user_id].get('title', '')
        content = self.params_cache[user_id].get('content', '')
        note = self.params_cache[user_id].get('note', '')
        # 将这些内容按照一定的格式整合到一起
        note_content = f"#sum4all\n{title}\n📒笔记：{note}\n{content}\n{link}"
        payload = {"content": note_content}
        # 将这个字典转换为JSON格式
        payload_json = json.dumps(payload)
        # 创建一个POST请求
        url = self.flomo_key
        headers = {'Content-Type': 'application/json'}
        # 发送这个POST请求
        response = requests.post(url, headers=headers, data=payload_json)
        reply = Reply()
        reply.type = ReplyType.TEXT
        if response.status_code == 200 and response.json()['code'] == 0:
            reply.content = f"已发送到{self.note_service}"        
        else:
            reply.content = "发送失败，错误码：" + str(response.status_code)
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS   
    def short_url(self, long_url):
        url = "https://short.fatwang2.com"
        payload = {
            "url": long_url
        }        
        headers = {'Content-Type': "application/json"}
        response = requests.request("POST", url, json=payload, headers=headers)
        if response.status_code == 200:
            res_data = response.json()
            # 直接从返回的 JSON 中获取短链接
            short_url = res_data.get('shorturl', None)  
            
            if short_url:
                return short_url
        return None
    def handle_url(self, content, e_context):
        logger.info('Handling Sum4All request...')
        # 根据sum_service的值选择API密钥和基础URL
        if self.url_sum_service == "openai":
            api_key = self.open_ai_api_key
            api_base = self.open_ai_api_base
            model = self.model
        elif self.url_sum_service == "sum4all":
            api_key = self.sum4all_key
            api_base = "https://pro.sum4all.site/v1"
            model = "sum4all"
        elif self.url_sum_service == "gemini":
            api_key = self.gemini_key
            model = "gemini"
            api_base = "https://gemini.sum4all.site/v1/models/gemini-pro:generateContent?key="
        else:
            logger.error(f"未知的sum_service配置: {self.url_sum_service}")
            return
        
        msg: ChatMessage = e_context["context"]["msg"]
        user_id = msg.from_user_id
        user_params = self.params_cache.get(user_id, {})
        isgroup = e_context["context"].get("isgroup", False)
        prompt = user_params.get('prompt', self.url_sum_prompt)
        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {api_key}'
        }
        payload = json.dumps({
            "link": content,
            "prompt": prompt,
            "model": model,
            "base": api_base
        })
        additional_content = ""  # 在 try 块之前初始化 additional_content

        try:
            logger.info('Sending request to LLM...')
            api_url = "https://ai.sum4all.site"
            response = requests.post(api_url, headers=headers, data=payload)
            response.raise_for_status()
            logger.info('Received response from LLM.')
            response_data = response.json()  # 解析响应的 JSON 数据
            if response_data.get("success"):
                content = response_data["content"].replace("\\n", "\n")  # 替换 \\n 为 \n
                self.params_cache[user_id]['content'] = content

                # 新增加的部分，用于解析 meta 数据
                meta = response_data.get("meta", {})  # 如果没有 meta 数据，则默认为空字典
                title = meta.get("og:title", "")  # 获取 og:title，如果没有则默认为空字符串
                self.params_cache[user_id]['title'] = title
                # 只有当 title 非空时，才加入到回复中
                if title:
                    additional_content += f"{title}\n\n"
                reply_content = additional_content + content  # 将内容加入回复
                
            else:
                reply_content = "Content not found or error in response"

        except requests.exceptions.RequestException as e:
            # 处理可能出现的错误
            logger.error(f"Error calling new combined api: {e}")
            reply_content = f"An error occurred"

        reply = Reply()
        reply.type = ReplyType.TEXT
        if isgroup or not self.note_enabled:
            reply.content = f"{remove_markdown(reply_content)}\n\n💬5min内输入{self.url_sum_qa_prefix}+问题，可继续追问"
        elif self.note_enabled:
            reply.content = f"{remove_markdown(reply_content)}\n\n💬5min内输入{self.url_sum_qa_prefix}+问题，可继续追问。\n\n📒输入{self.note_prefix}+笔记，可发送当前总结&笔记到{self.note_service}"
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    def handle_bibigpt(self, content, e_context):    
        headers = {
            'Content-Type': 'application/json'
        }
        payload_params = {
            "url": content,
            "includeDetail": False,
            "promptConfig": {
                "outputLanguage": self.outputLanguage
            }
        }

        payload = json.dumps(payload_params)           
        try:
            api_url = f"https://bibigpt.co/api/open/{self.bibigpt_key}"
            response = requests.request("POST",api_url, headers=headers, data=payload)
            response.raise_for_status()
            data = json.loads(response.text)
            summary_original = data.get('summary', 'Summary not available')
            html_url = data.get('htmlUrl', 'HTML URL not available')
            # 获取短链接
            short_url = self.short_url(html_url) 
            
            # 如果获取短链接失败，使用 html_url
            if short_url is None:
                short_url = html_url if html_url != 'HTML URL not available' else 'URL not available'
            
            # 移除 "##摘要"、"## 亮点" 和 "-"
            summary = summary_original.split("详细版（支持对话追问）")[0].replace("## 摘要\n", "📌总结：").replace("## 亮点\n", "").replace("- ", "")
        except requests.exceptions.RequestException as e:
            reply = f"An error occurred"

        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{summary}\n\n详细链接：{short_url}"

        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    def handle_opensum(self, content, e_context):
        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {self.opensum_key}'
        }
        payload = json.dumps({"link": content})
        try:
            api_url = "https://read.thinkwx.com/api/v1/article/summary"
            response = requests.request("POST",api_url, headers=headers, data=payload)
            response.raise_for_status()
            data = json.loads(response.text)
            summary_data = data.get('data', {})  # 获取data字段                
            summary_original = summary_data.get('summary', 'Summary not available')
            # 使用正则表达式提取URL
            url_pattern = r'https:\/\/[^\s]+'
            match = re.search(url_pattern, summary_original)
            html_url = match.group(0) if match else 'HTML URL not available'            
            # 获取短链接
            short_url = self.short_url(html_url) if match else html_url
            # 用于移除摘要中的URL及其后的所有内容
            url_pattern_remove = r'https:\/\/[^\s]+[\s\S]*'
            summary = re.sub(url_pattern_remove, '', summary_original).strip()        

        except requests.exceptions.RequestException as e:
            summary = f"An error occurred"
            short_url = 'URL not available'
        
        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{summary}\n\n详细链接：{short_url}"

        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS    
    def handle_search(self, content, e_context):
        # 根据sum_service的值选择API密钥和基础URL
        if self.search_sum_service == "openai":
            api_key = self.open_ai_api_key
            api_base = self.open_ai_api_base
            model = self.model
        elif self.search_sum_service == "sum4all":
            api_key = self.sum4all_key
            api_base = "https://pro.sum4all.site/v1"
            model = "sum4all"
        elif self.search_sum_service == "gemini":
            api_key = self.gemini_key
            model = "gemini"
            api_base = "https://gemini.sum4all.site/v1/models/gemini-pro:generateContent"

        else:
            logger.error(f"未知的search_service配置: {self.search_sum_service}")
            return
        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {api_key}'
        }
        content = content[len(self.search_sum_search_prefix):]
        payload = json.dumps({
            "ur": content,
            "prompt": self.search_sum_prompt,
            "model": model,
            "base": api_base
        })
        try:
            api_url = "https://ai.sum4all.site"
            response = requests.post(api_url, headers=headers, data=payload)
            response.raise_for_status()
            response_data = response.json()  # 解析响应的 JSON 数据
            if response_data.get("success"):
                content = response_data["content"].replace("\\n", "\n")  # 替换 \\n 为 \n
                reply_content = content  # 将内容加入回复

                # 解析 meta 数据
                meta = response_data.get("meta", {})  # 如果没有 meta 数据，则默认为空字典
                title = meta.get("og:title", "")  # 获取 og:title，如果没有则默认为空字符串
                og_url = meta.get("og:url", "")  # 获取 og:url，如果没有则默认为空字符串
                # 打印 title 和 og_url 以调试
                print("Title:", title)
                print("Original URL:", og_url)                
                # 只有当 title 和 url 非空时，才加入到回复中
                if title:
                    reply_content += f"\n\n参考文章：{title}"
                if og_url:
                    short_url = self.short_url(og_url)  # 获取短链接
                    reply_content += f"\n\n参考链接：{short_url}"                

            else:
                content = "Content not found or error in response"

        except requests.exceptions.RequestException as e:
            # 处理可能出现的错误
            logger.error(f"Error calling new combined api: {e}")
            reply_content = f"An error occurred"

        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{remove_markdown(reply_content)}"            
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    def handle_perplexity(self, content, e_context):

        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {self.perplexity_key}'
        }
        data = {
            "model": "sonar-small-online",
            "messages": [
                {"role": "system", "content": self.search_sum_prompt},
                {"role": "user", "content": content}
        ]
        }
        try:
            api_url = "https://api.perplexity.ai/chat/completions"
            response = requests.post(api_url, headers=headers, json=data)
            response.raise_for_status()
            # 处理响应数据
            response_data = response.json()
            # 这里可以根据你的需要处理响应数据
            # 解析 JSON 并获取 content
            if "choices" in response_data and len(response_data["choices"]) > 0:
                first_choice = response_data["choices"][0]
                if "message" in first_choice and "content" in first_choice["message"]:
                    content = first_choice["message"]["content"]
                else:
                    print("Content not found in the response")
            else:
                print("No choices available in the response")
        except requests.exceptions.RequestException as e:
            # 处理可能出现的错误
            logger.error(f"Error calling perplexity: {e}")
        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{remove_markdown(content)}"            
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    def get_help_text(self, **kwargs):
        help_text = "为你总结一切，发送链接、图片、文件，或直接搜索，一切都能为你总结\n"
        return help_text
    def handle_file(self, content, e_context):
        logger.info("handle_file: 向LLM发送内容总结请求")
        # 根据sum_service的值选择API密钥和基础URL
        if self.file_sum_service == "openai":
            api_key = self.open_ai_api_key
            api_base = self.open_ai_api_base
            model = self.model
        elif self.file_sum_service == "sum4all":
            api_key = self.sum4all_key
            api_base = "https://pro.sum4all.site/v1"
            model = "sum4all"
        elif self.file_sum_service == "gemini":
            api_key = self.gemini_key
            model = "gemini"
            api_base = "https://gemini.sum4all.site/v1/models/gemini-pro:generateContent"
        else:
            logger.error(f"未知的sum_service配置: {self.file_sum_service}")
            return
        msg: ChatMessage = e_context["context"]["msg"]
        user_id = msg.from_user_id
        user_params = self.params_cache.get(user_id, {})
        prompt = user_params.get('prompt', self.file_sum_prompt)
        if model == "gemini":
            headers = {
                'Content-Type': 'application/json',
                'x-goog-api-key': api_key
            }
            data = {
            "contents": [
                {"role": "user", "parts": [{"text": prompt}]},
                {"role": "model", "parts": [{"text": "okay"}]},
                {"role": "user", "parts": [{"text": content}]}
            ],
            "generationConfig": {
                "maxOutputTokens": 800
            }
            }
            api_url = api_base
        else:
            headers = {
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {api_key}'
            }
            data = {
                "model": model,
                "messages": [
                    {"role": "system", "content": prompt},
                    {"role": "user", "content": content}
                ]
            }
            api_url = f"{api_base}/chat/completions"
        try:
            response = requests.post(api_url, headers=headers, data=json.dumps(data))
            response.raise_for_status()
            response_data = response.json()
            
            # 解析 JSON 并获取 content
            if model == "gemini":
                if "candidates" in response_data and len(response_data["candidates"]) > 0:
                    first_candidate = response_data["candidates"][0]
                    if "content" in first_candidate:
                        if "parts" in first_candidate["content"] and len(first_candidate["content"]["parts"]) > 0:
                            response_content = first_candidate["content"]["parts"][0]["text"].strip()  # 获取响应内容
                            logger.info(f"Gemini API response content: {response_content}")  # 记录响应内容
                            reply_content = response_content.replace("\\n", "\n")  # 替换 \\n 为 \n
                        else:
                            logger.error("Parts not found in the Gemini API response content")
                            reply_content = "Parts not found in the Gemini API response content"
                    else:
                        logger.error("Content not found in the Gemini API response candidate")
                        reply_content = "Content not found in the Gemini API response candidate"
                else:
                    logger.error("No candidates available in the Gemini API response")
                    reply_content = "No candidates available in the Gemini API response"        
            else:
                if "choices" in response_data and len(response_data["choices"]) > 0:
                    first_choice = response_data["choices"][0]
                    if "message" in first_choice and "content" in first_choice["message"]:
                        response_content = first_choice["message"]["content"].strip()  # 获取响应内容
                        logger.info(f"LLM API response content")  # 记录响应内容
                        reply_content = response_content.replace("\\n", "\n")  # 替换 \\n 为 \n
                    else:
                        logger.error("Content not found in the response")
                        reply_content = "Content not found in the LLM API response"
                else:
                    logger.error("No choices available in the response")
                    reply_content = "No choices available in the LLM API response"

        except requests.exceptions.RequestException as e:
            logger.error(f"Error calling LLM API: {e}")
            reply_content = f"An error occurred while calling LLM API"

        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{remove_markdown(reply_content)}\n\n💬5min内输入{self.file_sum_qa_prefix}+问题，可继续追问" 
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    def read_pdf(self, file_path):
        logger.info(f"开始读取PDF文件：{file_path}")
        doc = fitz.open(file_path)
        content = ' '.join([page.get_text() for page in doc])
        logger.info(f"PDF文件读取完成：{file_path}")

        return content
    def read_word(self, file_path):
        doc = Document(file_path)
        return ' '.join([p.text for p in doc.paragraphs])
    def read_markdown(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            md_content = file.read()
            return markdown.markdown(md_content)
    def read_excel(self, file_path):
        workbook = load_workbook(file_path)
        content = ''
        for sheet in workbook:
            for row in sheet.iter_rows():
                content += ' '.join([str(cell.value) for cell in row])
                content += '\n'
        return content
    def read_txt(self, file_path):
        logger.debug(f"开始读取TXT文件: {file_path}")
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            logger.debug(f"TXT文件读取完成: {file_path}")
            logger.debug("TXT文件内容的前50个字符：")
            logger.debug(content[:50])  # 打印文件内容的前50个字符
            return content
        except Exception as e:
            logger.error(f"读取TXT文件时出错: {file_path}，错误信息: {str(e)}")
            return ""
    def read_csv(self, file_path):
        content = ''
        with open(file_path, 'r', encoding='utf-8') as csvfile:
            reader = csv.reader(csvfile)
            for row in reader:
                content += ' '.join(row) + '\n'
        return content
    def num_tokens_from_string(self, text):
        try:
            encoding = tiktoken.encoding_for_model("gpt-3.5-turbo")
        except KeyError:
            logger.debug(f"Warning: model not found. Using cl100k_base encoding.")
            encoding = tiktoken.get_encoding("cl100k_base")
        return len(encoding.encode(text))
    def read_html(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file, 'html.parser')
            return soup.get_text()
    def read_ppt(self, file_path):
        presentation = Presentation(file_path)
        content = ''
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + '\n'
        return content
    def split_text_chinese(self, text, overlap_tokens=500):
        tokens = jieba.cut(text)
        segments = []
        segment_text = ""
        for token in tokens:
            temp_segment_text = segment_text + token
            temp_segment_tokens_count = self.num_tokens_from_string(temp_segment_text)
            if temp_segment_tokens_count >= self.max_tokens:
                segments.append(segment_text)
                previous_segment_text = segment_text
                segment_text = previous_segment_text[-overlap_tokens:] + token if overlap_tokens > 0 else token
            else:
                segment_text = temp_segment_text

        if segment_text:
            segments.append(segment_text)
        logger.debug(f"分段文本: {segments}")
        return segments
    def extract_content(self, file_path):
        logger.info(f"extract_content: 提取文件内容，文件路径: {file_path}")

        file_extension = os.path.splitext(file_path)[1][1:].lower()
        logger.info(f"extract_content: 文件类型为 {file_extension}")

        file_type = EXTENSION_TO_TYPE.get(file_extension)

        if not file_type:
            logger.error(f"不支持的文件扩展名: {file_extension}")
            return None

        read_func = {
            'pdf': self.read_pdf,
            'docx': self.read_word,
            'md': self.read_markdown,
            'txt': self.read_txt,
            'excel': self.read_excel,
            'csv': self.read_csv,
            'html': self.read_html,
            'ppt': self.read_ppt
        }.get(file_type)

        if not read_func:
            logger.error(f"不支持的文件类型: {file_type}")
            return None
        logger.info("extract_content: 文件内容提取完成")
        return read_func(file_path)
    def encode_image_to_base64(self, image_path):
        # 打开图片
        img = Image.open(image_path)
        # 只有当图片的宽度大于1024像素时，才调整图片大小
        if img.width > 1024:
            img = img.resize((1024, int(img.height*1024/img.width)))
            # 将调整大小后的图片保存回原文件
            img.save(image_path)

        # 打开调整大小后的图片，读取并进行base64编码
        with open(image_path, "rb") as image_file:
            encoded = base64.b64encode(image_file.read()).decode('utf-8')
        return encoded
    # Function to handle OpenAI image processing
    def handle_openai_image(self, base64_image, e_context):
        logger.info("handle_openai_image_response: 解析OpenAI图像处理API的响应")
        msg: ChatMessage = e_context["context"]["msg"]
        user_id = msg.from_user_id
        user_params = self.params_cache.get(user_id, {})
        prompt = user_params.get('prompt', self.image_sum_prompt)

        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {self.open_ai_api_key}"
        }

        payload = {
            "model": "gpt-4-vision-preview",
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": prompt
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ],
            "max_tokens": 3000
        }

        try:
            response = requests.post(f"{self.open_ai_api_base}/chat/completions", headers=headers, json=payload)
            response.raise_for_status()  # 增加对HTTP错误的检查
            response_json = response.json()  # 定义response_json
            # 确保响应中有 'choices' 键并且至少有一个元素
            if "choices" in response_json and len(response_json["choices"]) > 0:
                first_choice = response_json["choices"][0]
                if "message" in first_choice and "content" in first_choice["message"]:
                    # 从响应中提取 'content'
                    response_content = first_choice["message"]["content"].strip()
                    logger.info("LLM API response content")  # 记录响应内容
                    reply_content = response_content
                else:
                    logger.error("Content not found in the response")
                    reply_content = "Content not found in the LLM API response"
            else:
                logger.error("No choices available in the response")
                reply_content = "No choices available in the LLM API response"
        except Exception as e:
            logger.error(f"Error processing LLM API response: {e}")
            reply_content = f"An error occurred while processing LLM API response"

        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{remove_markdown(reply_content)}\n\n💬5min内输入{self.image_sum_qa_prefix}+问题，可继续追问"  
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    def handle_gemini_image(self, base64_image, e_context):
        logger.info("handle_gemini_image: 解析Gemini图像处理API的响应")
        msg: ChatMessage = e_context["context"]["msg"]
        user_id = msg.from_user_id
        user_params = self.params_cache.get(user_id, {})
        prompt = user_params.get('prompt', self.image_sum_prompt)
        api_key = self.gemini_key
        payload = {
            "contents": [
                {
                    "parts": [
                        {"text": prompt},
                        {
                            "inline_data": {
                                "mime_type":"image/png",
                                "data": base64_image
                            }
                        }
                    ]
                }
            ],
            "generationConfig": {
                "maxOutputTokens": 800
            }
        }

        headers = {
            "Content-Type": "application/json",
            "x-goog-api-key": api_key
        }

        try:
            response = requests.post(f"https://gemini.sum4all.site/v1/models/gemini-pro-vision:generateContent", headers=headers, json=payload)
            response.raise_for_status()
            response_json = response.json()
            # 提取响应中的文本内容
            reply_content = response_json.get('candidates', [{}])[0].get('content', {}).get('parts', [{}])[0].get('text', 'No text found in the response')
        except Exception as e:
            reply_content = f"An error occurred while processing Gemini API response"

        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{remove_markdown(reply_content)}\n\n💬5min内输入{self.image_sum_qa_prefix}+问题，可继续追问"  
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    def handle_xunfei_image(self, base64_image, e_context):
        logger.info("handle_xunfei_image_response: 解析讯飞图像处理API的响应")
        msg: ChatMessage = e_context["context"]["msg"]
        user_id = msg.from_user_id
        user_params = self.params_cache.get(user_id, {})
        prompt = user_params.get('prompt', self.image_sum_prompt)

        headers = {
            "Content-Type": "application/json",
            "X_API_KEY": self.xunfei_api_key,
            "X_APP_ID": self.xunfei_app_id,
            "X_API_SECRET": self.xunfei_api_secret
        }

        payload = {
            "model": "spark-chat-vision",
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        },
                        {
                            "type": "text",
                            "text": prompt
                        }
                    ]
                }
            ],
            "max_tokens": 3000
        }

        try:
            response = requests.post("https://spark.sum4all.site/v1/chat/completions", headers=headers, json=payload)
            response.raise_for_status()  # 增加对HTTP错误的检查
            response_json = response.json()  # 定义response_json
            # 确保响应中有 'choices' 键并且至少有一个元素
            if "choices" in response_json and len(response_json["choices"]) > 0:
                first_choice = response_json["choices"][0]
                if "message" in first_choice and "content" in first_choice["message"]:
                    # 从响应中提取 'content'
                    response_content = first_choice["message"]["content"].strip()
                    logger.info("LLM API response content")  # 记录响应内容
                    reply_content = response_content
                else:
                    logger.error("Content not found in the response")
                    reply_content = "Content not found in the LLM API response"
            else:
                logger.error("No choices available in the response")
                reply_content = "No choices available in the LLM API response"
        except Exception as e:
            logger.error(f"Error processing LLM API response: {e}")
            reply_content = f"An error occurred while processing LLM API response"

        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{remove_markdown(reply_content)}\n\n💬5min内输入{self.image_sum_qa_prefix}+问题，可继续追问"  
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    
def remove_markdown(text):
    # 替换Markdown的粗体标记
    text = text.replace("**", "")
    # 替换Markdown的标题标记
    text = text.replace("### ", "").replace("## ", "").replace("# ", "")
    return text